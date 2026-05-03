#!/usr/bin/env python3
"""
KI-Vortrag PowerPoint Generator
Konvertiert 02-KI-Praesentation.md → KI-Vortrag.pptx  (eigenes Design)
                                   → KI-Vortrag-Vorlage.pptx  (basierend auf Vorlage.pptx)

Design: Modern Flat  |  Navy #1E3A5F  |  Amber #F5A623  |  16:9
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Farben ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_DARK   = RGBColor(0x16, 0x2D, 0x4B)
AMBER       = RGBColor(0xF5, 0xA6, 0x23)
AMBER_L     = RGBColor(0xFF, 0xC9, 0x5A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x2C, 0x3E, 0x50)
GREY        = RGBColor(0x7F, 0x8C, 0x8D)
LGREY       = RGBColor(0xF0, 0xF4, 0xF8)
BGREY       = RGBColor(0xB0, 0xC4, 0xDE)
TALT        = RGBColor(0xEE, 0xF2, 0xF7)
BLUE_MUTED  = RGBColor(0x60, 0x80, 0xA0)

# ── Maße (16:9) ───────────────────────────────────────────────────────────────
SW = Inches(13.333)   # slide width
SH = Inches(7.500)    # slide height
MG = Inches(0.40)     # margin
HH = Inches(1.05)     # header bar height
IW = Inches(4.40)     # image placeholder width
CT = MG + HH + Inches(0.12)   # content area top
CH = SH - CT - MG             # content area height
CW = SW - 2 * MG              # full content width


# ── Hilfsfunktionen ──────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None, line_pt=1.0):
    """Rechteck-Shape hinzufügen."""
    s = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def parse_md(text: str):
    """**bold**, *italic*, `code` → [(text, bold, italic), ...]"""
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)   # Links entfernen
    result = []
    pat = re.compile(r'\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`|([^*`]+)', re.DOTALL)
    for m in pat.finditer(text):
        if m.group(1):   result.append((m.group(1), True,  False))
        elif m.group(2): result.append((m.group(2), False, True))
        elif m.group(3): result.append((m.group(3), False, False))
        elif m.group(4): result.append((m.group(4), False, False))
    return result or [(text, False, False)]


def add_run(para, text, size=16, bold=False, italic=False, color=DARK):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def calc_font(n: int) -> int:
    """Schriftgröße basierend auf Anzahl Bullets."""
    if n <= 4:  return 17
    if n <= 6:  return 15
    if n <= 9:  return 13
    return 11


def add_header(slide, title: str, slide_num=None):
    """Navy-Header-Leiste mit Amber-Akzent und weißem Titel."""
    add_rect(slide, MG, MG, SW - 2*MG, HH, fill=NAVY)
    add_rect(slide, MG, MG, Inches(0.07), HH, fill=AMBER)

    tb = slide.shapes.add_textbox(
        MG + Inches(0.14), MG + Inches(0.08),
        SW - 2*MG - Inches(0.14), HH - Inches(0.08))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_run(p, title, size=24, bold=True, color=WHITE)

    # Untere Amber-Linie
    add_rect(slide, MG, SH - MG - Inches(0.05), SW - 2*MG, Inches(0.05), fill=AMBER)

    # Foliennummer
    if slide_num is not None:
        nb = slide.shapes.add_textbox(
            SW - Inches(1.3), SH - MG - Inches(0.26),
            Inches(1.1), Inches(0.22))
        p2 = nb.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        add_run(p2, str(slide_num), size=10, color=GREY)


def add_img_ph(slide, prompt: str):
    """Bildplatzhalter-Box rechts."""
    l = SW - MG - IW
    add_rect(slide, l, CT, IW, CH, fill=LGREY, line=BGREY, line_pt=1.5)

    tb = slide.shapes.add_textbox(
        l + Inches(0.12), CT + Inches(0.1),
        IW - Inches(0.24), CH - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_run(p, "🖼  Bildplatzhalter", size=10, bold=True, color=NAVY)

    if prompt:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        add_run(p2, "\n" + prompt, size=8, italic=True, color=GREY)


def add_bullets(slide, bullets: list, right_margin=0):
    """Bullet-Liste in den Content-Bereich einfügen."""
    n = len([b for b in bullets if not b.get('section_header')])
    fs = calc_font(n)
    cw = CW - right_margin

    tb = slide.shapes.add_textbox(MG, CT, cw, CH)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        if b.get('section_header'):
            p.space_before = Pt(7)
            p.space_after  = Pt(2)
            for txt, bold, ital in parse_md(b['text']):
                if txt:
                    add_run(p, txt, size=fs, bold=True, color=NAVY)
            continue

        p.space_before = Pt(3)
        p.space_after  = Pt(1)
        indent = b.get('indent', 0)
        num    = b.get('num')

        if num is not None:
            add_run(p, f"  {num}.  ", size=fs, bold=True, color=AMBER)
        elif indent:
            add_run(p, "       –  ", size=fs - 1, bold=True, color=AMBER)
        else:
            add_run(p, "  •  ", size=fs, bold=True, color=AMBER)

        for txt, bold, ital in parse_md(b['text']):
            if txt:
                add_run(p, txt, size=fs - (1 if indent else 0),
                        bold=bold, italic=ital, color=DARK)


def add_table(slide, rows: list, left, top, width, height):
    """Formatierte Tabelle einfügen."""
    if not rows:
        return
    nr, nc = len(rows), len(rows[0])
    rh = min(Inches(0.50), height / nr)

    tbl = slide.shapes.add_table(
        nr, nc, int(left), int(top), int(width), int(rh * nr)).table

    # Header-Zeile
    for j, txt in enumerate(rows[0]):
        c = tbl.cell(0, j)
        c.text = ''.join(s[0] for s in parse_md(txt))
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        for para in c.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = WHITE

    # Daten-Zeilen
    for i, row in enumerate(rows[1:], 1):
        bg = TALT if i % 2 == 0 else WHITE
        for j, txt in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ''.join(s[0] for s in parse_md(txt))
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.bold = (j == 0)
                    run.font.color.rgb = DARK


def set_notes(slide, text: str):
    if text and text.strip():
        slide.notes_slide.notes_text_frame.text = text.strip()


# ── Markdown Parser ───────────────────────────────────────────────────────────

def parse_markdown(path: str) -> list:
    """
    Parst die Präsentation in eine Liste von Slide-Dicts.
    Behandelt Folie 11 (Orphan-Block nach ---) korrekt.
    """
    text = Path(path).read_text('utf-8')

    # Blöcke trennen, Orphan-Blöcke mit vorherigem zusammenführen
    raw = re.split(r'\n---\n', text)
    blocks = []
    for blk in raw:
        blk = blk.strip()
        if not blk:
            continue
        first_line = blk.split('\n')[0].strip()
        is_known = (
            first_line.startswith('## Folie') or
            first_line.startswith('## Teil') or
            first_line.startswith('# ') or
            '## Anhang' in first_line
        )
        if is_known or not blocks:
            blocks.append(blk)
        else:
            blocks[-1] = blocks[-1] + '\n\n' + blk   # z.B. Folie 11

    slides = []

    for block in blocks:
        lines = block.split('\n')
        if not lines:
            continue
        first = lines[0].strip()

        # ── Abschnitts-Trennfolie (Teil N) ────────────────────────────────
        m = re.match(r'^## Teil (\d+)[:\.]?\s*(.+)$', first)
        if m:
            slides.append({
                'type':  'section',
                'num':   int(m.group(1)),
                'title': m.group(2).strip(),
                'notes': '',
            })
            continue

        # ── Glossar ────────────────────────────────────────────────────────
        if first.startswith('## Anhang'):
            table_buf = []
            for line in lines:
                if line.startswith('|'):
                    if re.match(r'^\|[-| :]+\|$', line.strip()):
                        continue
                    cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    table_buf.append(cells)
            slides.append({
                'type':       'glossar',
                'title':      'Anhang: Glossar',
                'table_data': table_buf,
                'notes':      '',
                'image_prompt': '',
                'bullets':    [],
            })
            continue

        # ── Reguläre Folie ─────────────────────────────────────────────────
        m = re.match(r'^## Folie (\d+)\s+(?:–|-)\s+(.+)$', first)
        if not m:
            continue

        num   = int(m.group(1))
        title = m.group(2).strip()

        sd = {
            'type':         'content',
            'num':          num,
            'title':        title,
            'bullets':      [],
            'table_data':   [],
            'image_prompt': '',
            'notes':        '',
            'main_text':    title,          # für Titelfolie
            'subtitle_text': '',
        }

        in_notes = False
        in_img   = False
        notes_buf = []
        table_buf = []

        for line in lines[1:]:
            # Leere Blockquote-Zeile
            if line.strip() == '>':
                if in_notes:
                    notes_buf.append('')
                continue

            # Blockquote-Zeile
            if line.startswith('> '):
                content = line[2:].strip()

                if '**Sprecher-Notizen:**' in content:
                    in_notes = True;  in_img = False
                elif '**Bild-Prompt:**' in content:
                    in_img = True;    in_notes = False
                elif content.startswith('*Quelle:') or content.startswith('*Quelle'):
                    in_notes = False; in_img = False
                elif in_img and content.startswith('`') and content.endswith('`'):
                    sd['image_prompt'] = content[1:-1]
                    in_img = False
                elif in_notes and content:
                    clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
                    clean = re.sub(r'\*([^*]+)\*', r'\1', clean)
                    notes_buf.append(clean)
                continue

            # Tabellen-Zeile
            if line.startswith('|'):
                if re.match(r'^\|[-| :]+\|$', line.strip()):
                    continue
                cells = [c.strip() for c in line.strip().strip('|').split('|')]
                table_buf.append(cells)
                continue

            # Nummerierte Liste
            nm = re.match(r'^(\d+)\.\s+(.+)$', line)
            if nm:
                sd['bullets'].append({
                    'text':           nm.group(2).strip(),
                    'num':            int(nm.group(1)),
                    'section_header': False,
                    'indent':         0,
                })
                continue

            # Bullet
            if line.startswith('- '):
                sd['bullets'].append({
                    'text':           line[2:].strip(),
                    'num':            None,
                    'section_header': False,
                    'indent':         0,
                })
                continue

            # Inline-Abschnitts-Header innerhalb einer Folie: **Text**
            sh = re.match(r'^\*\*(.+)\*\*$', line.strip())
            if sh:
                sd['bullets'].append({
                    'text':           sh.group(1),
                    'num':            None,
                    'section_header': True,
                    'indent':         0,
                })
                continue

            # Titelfolie: **Titel** und *Untertitel*
            if num == 1:
                bt = re.match(r'^\*\*(.+)\*\*$', line.strip())
                if bt:
                    sd['main_text'] = bt.group(1)
                    continue
                it = re.match(r'^\*([^*].+)\*$', line.strip())
                if it:
                    sd['subtitle_text'] = it.group(1)
                    continue

        if table_buf:
            sd['table_data'] = table_buf
        sd['notes'] = '\n'.join(notes_buf).strip()

        # Slide-Typ bestimmen
        if num == 1:
            sd['type'] = 'title'
        elif num == 36:
            sd['type'] = 'quote'
        elif sd['table_data'] and not sd['bullets']:
            sd['type'] = 'table'
        elif sd['table_data'] and sd['bullets']:
            sd['type'] = 'mixed'

        slides.append(sd)

    return slides


# ── Slide Builder ─────────────────────────────────────────────────────────────

def build_title(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Rechtes Panel (dunkler)
    add_rect(slide, SW - Inches(5.2), 0, Inches(5.2), SH, fill=NAVY_DARK)
    # Untere Amber-Leiste
    add_rect(slide, 0, SH - Inches(0.45), SW, Inches(0.45), fill=AMBER)
    # Amber-Vertikalakzent
    add_rect(slide, MG, Inches(1.5), Inches(0.07), Inches(3.5), fill=AMBER)

    # Haupttitel
    tb = slide.shapes.add_textbox(
        MG + Inches(0.22), Inches(1.6), Inches(7.6), Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_run(p, sd.get('main_text', 'Künstliche Intelligenz'),
            size=46, bold=True, color=WHITE)

    # Untertitel
    tb2 = slide.shapes.add_textbox(
        MG + Inches(0.22), Inches(3.6), Inches(7.6), Inches(0.9))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    add_run(p2,
            sd.get('subtitle_text', 'Verstehen. Nutzen. Verantwortungsvoll handeln.'),
            size=21, italic=True, color=RGBColor(0xA0, 0xC0, 0xE0))

    # Dauer
    tb3 = slide.shapes.add_textbox(
        MG + Inches(0.22), Inches(4.6), Inches(5), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]
    add_run(p3, "Einführung · ca. 60 Minuten", size=13, color=GREY)

    # Bildplatzhalter im rechten Panel
    if sd.get('image_prompt'):
        tb4 = slide.shapes.add_textbox(
            SW - Inches(4.9), Inches(0.9), Inches(4.4), Inches(5.5))
        tf4 = tb4.text_frame; tf4.word_wrap = True
        p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
        add_run(p4, "🖼  Bildplatzhalter", size=10, bold=True, color=BLUE_MUTED)
        p5 = tf4.add_paragraph(); p5.alignment = PP_ALIGN.LEFT
        add_run(p5, "\n" + sd['image_prompt'], size=8, italic=True, color=BLUE_MUTED)

    set_notes(slide, sd.get('notes', ''))
    return slide


def build_section(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # Amber-Trennlinie
    add_rect(slide, Inches(1.5), SH / 2 + Inches(0.38),
             Inches(10.3), Inches(0.07), fill=AMBER)

    # Teil N
    tb1 = slide.shapes.add_textbox(0, SH / 2 - Inches(1.15), SW, Inches(0.65))
    p1 = tb1.text_frame.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    add_run(p1, f"Teil {sd['num']}", size=17, bold=True, color=AMBER_L)

    # Titel
    tb2 = slide.shapes.add_textbox(
        Inches(0.8), SH / 2 - Inches(0.52), SW - Inches(1.6), Inches(1.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    add_run(p2, sd['title'], size=34, bold=True, color=WHITE)

    return slide


def build_content(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    has_img = bool(sd.get('image_prompt'))
    rm = IW + Inches(0.25) if has_img else 0

    add_header(slide, sd['title'], sd['num'])

    if sd['bullets']:
        add_bullets(slide, sd['bullets'], right_margin=rm)

    if has_img:
        add_img_ph(slide, sd['image_prompt'])

    set_notes(slide, sd.get('notes', ''))
    return slide


def build_table_slide(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    has_img = bool(sd.get('image_prompt'))
    rm = IW + Inches(0.25) if has_img else 0

    add_header(slide, sd['title'], sd.get('num'))

    if sd.get('table_data'):
        add_table(slide, sd['table_data'], MG, CT, CW - rm, CH)

    if has_img:
        add_img_ph(slide, sd['image_prompt'])

    set_notes(slide, sd.get('notes', ''))
    return slide


def build_mixed(prs, sd):
    """Folie mit Bullets UND Tabelle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    has_img = bool(sd.get('image_prompt'))
    rm = IW + Inches(0.25) if has_img else 0
    cw = CW - rm

    add_header(slide, sd['title'], sd['num'])

    # Obere Hälfte: Bullets
    bh = CH * 0.50
    tb = slide.shapes.add_textbox(MG, CT, cw, bh)
    tf = tb.text_frame; tf.word_wrap = True
    n = len([b for b in sd['bullets'] if not b.get('section_header')])
    fs = calc_font(n)

    for i, b in enumerate(sd['bullets']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5 if b.get('section_header') else 3)
        p.alignment = PP_ALIGN.LEFT
        if b.get('section_header'):
            for txt, bold, ital in parse_md(b['text']):
                if txt: add_run(p, txt, size=fs, bold=True, color=NAVY)
        else:
            add_run(p, "  •  ", size=fs, bold=True, color=AMBER)
            for txt, bold, ital in parse_md(b['text']):
                if txt: add_run(p, txt, size=fs, bold=bold, italic=ital, color=DARK)

    # Untere Hälfte: Tabelle
    if sd.get('table_data'):
        table_top = CT + bh + Inches(0.08)
        table_h   = CH - bh - Inches(0.08)
        add_table(slide, sd['table_data'], MG, table_top, cw, table_h)

    if has_img:
        add_img_ph(slide, sd['image_prompt'])

    set_notes(slide, sd.get('notes', ''))
    return slide


def build_quote(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_rect(slide, 0, SH - Inches(0.45), SW, Inches(0.45), fill=AMBER)
    add_rect(slide, SW - Inches(5.0), 0, Inches(5.0), SH, fill=NAVY_DARK)

    # Zitat-Text
    quotes = [
        "„KI ist nicht die Zukunft.",
        "KI ist die Gegenwart.",
        "Die Frage ist nur:",
        "Wie gestalten wir sie?"
    ]
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(7.4), Inches(4.6))
    tf = tb.text_frame; tf.word_wrap = True

    for i, line in enumerate(quotes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        add_run(p, line, size=30, bold=True, color=WHITE)

    # Bildplatzhalter rechts
    if sd.get('image_prompt'):
        add_img_ph(slide, sd['image_prompt'])

    set_notes(slide, sd.get('notes', ''))
    return slide


def build_glossar(prs, sd):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)

    add_rect(slide, MG, MG, SW - 2*MG, HH, fill=NAVY)
    add_rect(slide, MG, MG, Inches(0.07), HH, fill=AMBER)
    add_rect(slide, MG, SH - MG - Inches(0.05), SW - 2*MG, Inches(0.05), fill=AMBER)

    tb = slide.shapes.add_textbox(
        MG + Inches(0.14), MG + Inches(0.08),
        SW - 2*MG - Inches(0.14), HH - Inches(0.08))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    add_run(p, "Anhang: Glossar", size=24, bold=True, color=WHITE)

    if sd.get('table_data'):
        add_table(slide, sd['table_data'], MG, CT, CW, CH)

    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

BUILDERS = {
    'title':   build_title,
    'section': build_section,
    'content': build_content,
    'table':   build_table_slide,
    'mixed':   build_mixed,
    'quote':   build_quote,
    'glossar': build_glossar,
}


def main():
    md_path  = Path(__file__).parent / '02-KI-Praesentation.md'
    out_path = Path(__file__).parent / 'KI-Vortrag.pptx'

    print(f"Lese: {md_path}")
    slides_data = parse_markdown(str(md_path))

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print(f"\n{len(slides_data)} Elemente gefunden:\n")
    for sd in slides_data:
        t = sd.get('type', 'content')
        label = sd.get('title', f"Teil {sd.get('num', '?')}")
        n = sd.get('num', '')
        num_str = f"[{n:>3}]" if n else "     "
        print(f"  {num_str} [{t:8s}]  {label[:60]}")
        BUILDERS.get(t, build_content)(prs, sd)

    prs.save(str(out_path))
    print(f"\n✓ Gespeichert: {out_path}")
    print(f"  {len(prs.slides)} Folien erstellt")


def main_vorlage():
    """Zweite Variante: verwendet Vorlage.pptx als Basis (Slide-Master + Layouts)."""
    from pptx.util import Emu
    from copy import deepcopy
    import lxml.etree as etree

    md_path     = Path(__file__).parent / '02-KI-Praesentation.md'
    tmpl_path   = Path(__file__).parent / 'Vorlage.pptx'
    out_path    = Path(__file__).parent / 'KI-Vortrag-Vorlage.pptx'

    print(f"Lese: {md_path}")
    print(f"Vorlage: {tmpl_path}")
    slides_data = parse_markdown(str(md_path))

    # ── Vorlage laden ────────────────────────────────────────────────────────
    prs = Presentation(str(tmpl_path))

    # Demo-Folien aus der Vorlage sauber entfernen (inkl. Relationship-Drop)
    slides_list = prs.slides._sldIdLst
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    for sld_id in list(slides_list):
        rId = sld_id.get(f'{{{NS_R}}}id')
        if rId:
            prs.part.drop_rel(rId)
        slides_list.remove(sld_id)

    # Layout-Zuordnung (nach Name, Fallback: Leer)
    layout_map = {l.name: l for l in prs.slide_layouts}
    L_TITLE    = layout_map.get('Titelfolie',           prs.slide_layouts[0])
    L_CONTENT  = layout_map.get('Titel und Inhalt',     prs.slide_layouts[1])
    L_SECTION  = layout_map.get('Abschnitts-\nüberschrift', prs.slide_layouts[2])
    L_BLANK    = layout_map.get('Leer',                 prs.slide_layouts[6])
    L_QUOTE    = layout_map.get('Zitat mit Beschriftung', layout_map.get('Leer', prs.slide_layouts[6]))

    def _add_slide(layout):
        """Folie mit gegebenem Layout hinzufügen."""
        slide_layout = layout
        rId = prs.slides._sldIdLst.getparent().part.relate_to(
            slide_layout.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout')
        # python-pptx interner Weg
        from pptx.parts.slide import SlidePart
        from pptx.oxml.ns import qn
        slide_part = prs.slides._sldIdLst.getparent().part.add_slide(slide_layout)
        slide = slide_part.slide
        return slide

    # Einfacherer Weg: Presentation.slides.add_slide
    def add_slide(layout):
        return prs.slides.add_slide(layout)

    def _get_ph(slide, idx):
        """Placeholder nach idx holen, None wenn nicht vorhanden."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_title(slide, title_text, size=28, bold=True, color=DARK):
        ph = _get_ph(slide, 0)
        if ph:
            ph.text = title_text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size)
                    run.font.bold = bold
                    run.font.color.rgb = color

    def _fill_content_ph(slide, bullets, tables, image_prompt):
        """Inhalt-Placeholder (idx=1) mit Bullets oder Tabelle füllen."""
        ph = _get_ph(slide, 1)

        # Wenn kein Content-Placeholder vorhanden: Textbox hinzufügen
        if ph is None:
            tb = slide.shapes.add_textbox(
                Inches(0.75), Inches(2.34), Inches(11.08), Inches(3.99))
            tf = tb.text_frame
            tf.word_wrap = True
            _write_bullets(tf, bullets, image_prompt)
            return

        tf = ph.text_frame
        tf.word_wrap = True
        tf.clear()
        _write_bullets(tf, bullets, image_prompt)

    def _write_bullets(tf, bullets, image_prompt):
        """Bullets in ein TextFrame schreiben."""
        n  = len([b for b in bullets if not b.get('section_header')])
        fs = calc_font(n)
        first = True

        for b in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4 if b.get('section_header') else 2)
            p.alignment    = PP_ALIGN.LEFT

            if b.get('section_header'):
                p.level = 0
                for txt, bold, ital in parse_md(b['text']):
                    if txt:
                        r = p.add_run()
                        r.text = txt
                        r.font.size  = Pt(fs)
                        r.font.bold  = True
                        r.font.color.rgb = NAVY
            else:
                p.level = 1 if b.get('indent') else 0
                num = b.get('num')
                prefix = f"{num}.  " if num is not None else "•  "
                r0 = p.add_run()
                r0.text = prefix
                r0.font.size  = Pt(fs)
                r0.font.bold  = True
                r0.font.color.rgb = AMBER

                for txt, bold, ital in parse_md(b['text']):
                    if txt:
                        r = p.add_run()
                        r.text       = txt
                        r.font.size  = Pt(fs - (1 if b.get('indent') else 0))
                        r.font.bold  = bold
                        r.font.italic = ital
                        r.font.color.rgb = DARK

        # Bildplatzhalter als letzten Absatz anhängen
        if image_prompt:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(10)
            r = p2.add_run()
            r.text = f"🖼  Bildplatzhalter: {image_prompt[:120]}..."
            r.font.size    = Pt(9)
            r.font.italic  = True
            r.font.color.rgb = GREY

    def _add_table_to_slide(slide, table_data):
        """Tabelle direkt auf die Folie legen (unter Titel-Bereich)."""
        if not table_data:
            return
        nr, nc = len(table_data), len(table_data[0])
        l = Inches(0.75); t = Inches(2.34)
        w = Inches(11.08); h = Inches(min(0.5 * nr, 3.9))

        tbl = slide.shapes.add_table(nr, nc, int(l), int(t), int(w), int(h)).table

        for j, cell_txt in enumerate(table_data[0]):
            c = tbl.cell(0, j)
            c.text = ''.join(s[0] for s in parse_md(cell_txt))
            c.fill.solid(); c.fill.fore_color.rgb = NAVY
            for para in c.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(12); run.font.bold = True
                    run.font.color.rgb = WHITE

        for i, row in enumerate(table_data[1:], 1):
            bg = TALT if i % 2 == 0 else WHITE
            for j, cell_txt in enumerate(row):
                c = tbl.cell(i, j)
                c.text = ''.join(s[0] for s in parse_md(cell_txt))
                c.fill.solid(); c.fill.fore_color.rgb = bg
                for para in c.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.size = Pt(11); run.font.bold = (j == 0)
                        run.font.color.rgb = DARK

    # ── Folie-Builder (Vorlage-Version) ─────────────────────────────────────

    def vb_title(sd):
        slide = add_slide(L_TITLE)
        ph0 = _get_ph(slide, 0)
        ph1 = _get_ph(slide, 1)
        if ph0:
            ph0.text = sd.get('main_text', 'Künstliche Intelligenz')
            for para in ph0.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold  = True
                    run.font.color.rgb = NAVY
        if ph1:
            ph1.text = sd.get('subtitle_text', 'Verstehen. Nutzen. Verantwortungsvoll handeln.')
        set_notes(slide, sd.get('notes', ''))
        return slide

    def vb_section(sd):
        slide = add_slide(L_SECTION)
        ph0 = _get_ph(slide, 0)
        if ph0:
            ph0.text = sd['title']
            for para in ph0.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size  = Pt(32)
                    run.font.bold  = True
        ph1 = _get_ph(slide, 1)
        if ph1:
            ph1.text = f"Teil {sd['num']}"
            for para in ph1.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)
        return slide

    def vb_content(sd):
        slide = add_slide(L_CONTENT)
        _set_title(slide, sd['title'])
        _fill_content_ph(slide, sd.get('bullets', []),
                         sd.get('table_data', []), sd.get('image_prompt', ''))
        set_notes(slide, sd.get('notes', ''))
        return slide

    def vb_table(sd):
        slide = add_slide(L_CONTENT)
        _set_title(slide, sd['title'])
        # Placeholder leeren, dann Tabelle als Shape hinzufügen
        ph = _get_ph(slide, 1)
        if ph:
            from pptx.oxml.ns import qn
            sp = ph._element
            sp.getparent().remove(sp)
        _add_table_to_slide(slide, sd.get('table_data', []))
        if sd.get('image_prompt'):
            tb = slide.shapes.add_textbox(
                Inches(0.75), Inches(6.3), Inches(11.08), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = f"🖼  {sd['image_prompt'][:130]}..."
            r.font.size = Pt(8); r.font.italic = True
            r.font.color.rgb = GREY
        set_notes(slide, sd.get('notes', ''))
        return slide

    def vb_mixed(sd):
        slide = add_slide(L_CONTENT)
        _set_title(slide, sd['title'])
        ph = _get_ph(slide, 1)
        if ph:
            from pptx.oxml.ns import qn
            sp = ph._element
            sp.getparent().remove(sp)

        # Bullets oben
        bh = Inches(1.8)
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.34), Inches(11.08), bh)
        tf = tb.text_frame; tf.word_wrap = True
        _write_bullets(tf, sd.get('bullets', []), '')

        # Tabelle unten
        rows = sd.get('table_data', [])
        if rows:
            nr, nc = len(rows), len(rows[0])
            tt = Inches(2.34) + bh + Inches(0.1)
            th = Inches(7.5) - tt - Inches(0.5)
            tbl = slide.shapes.add_table(
                nr, nc, int(Inches(0.75)), int(tt),
                int(Inches(11.08)), int(th)).table
            for j, ct in enumerate(rows[0]):
                c = tbl.cell(0, j); c.text = ''.join(s[0] for s in parse_md(ct))
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
                for para in c.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(11); run.font.bold = True
                        run.font.color.rgb = WHITE
            for i, row in enumerate(rows[1:], 1):
                bg = TALT if i % 2 == 0 else WHITE
                for j, ct in enumerate(row):
                    c = tbl.cell(i, j); c.text = ''.join(s[0] for s in parse_md(ct))
                    c.fill.solid(); c.fill.fore_color.rgb = bg
                    for para in c.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(10); run.font.bold = (j == 0)
                            run.font.color.rgb = DARK

        set_notes(slide, sd.get('notes', ''))
        return slide

    def vb_quote(sd):
        slide = add_slide(L_QUOTE)
        # Alle vorhandenen Placeholder füllen oder Textbox nutzen
        ph0 = _get_ph(slide, 0)
        quote = (
            '„KI ist nicht die Zukunft.\n'
            'KI ist die Gegenwart.\n'
            'Die Frage ist nur:\n'
            'Wie gestalten wir sie?"'
        )
        if ph0:
            ph0.text = quote
            for para in ph0.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(26); run.font.bold = True
        else:
            tb = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.6), Inches(11.0), Inches(4.5))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = quote
            r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = DARK
        if sd.get('image_prompt'):
            tb2 = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.0), Inches(11.0), Inches(0.7))
            p2 = tb2.text_frame.paragraphs[0]
            r2 = p2.add_run()
            r2.text = f"🖼  {sd['image_prompt'][:130]}..."
            r2.font.size = Pt(8); r2.font.italic = True; r2.font.color.rgb = GREY
        set_notes(slide, sd.get('notes', ''))
        return slide

    def vb_glossar(sd):
        slide = add_slide(L_CONTENT)
        _set_title(slide, 'Anhang: Glossar')
        ph = _get_ph(slide, 1)
        if ph:
            from pptx.oxml.ns import qn
            ph._element.getparent().remove(ph._element)
        _add_table_to_slide(slide, sd.get('table_data', []))
        return slide

    VBUILDERS = {
        'title':   vb_title,
        'section': vb_section,
        'content': vb_content,
        'table':   vb_table,
        'mixed':   vb_mixed,
        'quote':   vb_quote,
        'glossar': vb_glossar,
    }

    # ── Folien erzeugen ──────────────────────────────────────────────────────
    print(f"\n{len(slides_data)} Elemente gefunden:\n")
    for sd in slides_data:
        t = sd.get('type', 'content')
        label = sd.get('title', f"Teil {sd.get('num', '?')}")
        n = sd.get('num', '')
        num_str = f"[{n:>3}]" if n else "     "
        print(f"  {num_str} [{t:8s}]  {label[:60]}")
        VBUILDERS.get(t, vb_content)(sd)

    prs.save(str(out_path))
    print(f"\n✓ Gespeichert: {out_path}")
    print(f"  {len(prs.slides)} Folien erstellt")


if __name__ == '__main__':
    import sys
    if '--vorlage' in sys.argv:
        main_vorlage()
    else:
        main()
        main_vorlage()
