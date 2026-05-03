#!/usr/bin/env python3
"""
KI-VHS-Vortrag Generator
Template : KI_VHS_60min_im_Stil_Vorlage01.pptx
Output   : 03_KI_VHS_Vortrag.pptx
Grafiken : Grafiken/NN_Folientitel.txt  (Bild-Prompts im Comic-Stil)

Bildstil für alle Grafiken:
  flat vector comic illustration, bold outlines, vibrant colors,
  clean white background, friendly and modern, no text in image,
  professional presentation quality, 16:9 aspect ratio
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re, textwrap

# ── Pfade ─────────────────────────────────────────────────────────────────────
BASE    = Path(__file__).parent
TMPL    = BASE / 'KI_VHS_60min_im_Stil_Vorlage01.pptx'
MD      = BASE / '02-KI-Praesentation.md'
OUT     = BASE / '03_KI_VHS_Vortrag.pptx'
GRAFDIR = BASE / 'Grafiken'
GRAFDIR.mkdir(exist_ok=True)

# ── Globaler Bildstil (an jeden Prompt anhängen) ──────────────────────────────
IMG_STYLE = (
    "flat vector comic illustration, bold outlines, vibrant saturated colors, "
    "clean white background, friendly modern style, no text in image, "
    "professional presentation quality, 16:9 aspect ratio"
)

# ── Hilfsfunktionen ───────────────────────────────────────────────────────────
def slugify(text: str) -> str:
    """Titeltext → sicherer Dateiname."""
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'\s+', '_', text.strip())
    return text[:50]


def parse_md(text: str):
    """**bold**, *italic*, `code` → [(text, bold, italic)]"""
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    result = []
    pat = re.compile(r'\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`|([^*`]+)', re.DOTALL)
    for m in pat.finditer(text):
        if m.group(1):   result.append((m.group(1), True,  False))
        elif m.group(2): result.append((m.group(2), False, True))
        elif m.group(3): result.append((m.group(3), False, False))
        elif m.group(4) and m.group(4).strip(): result.append((m.group(4), False, False))
    return result or [(text, False, False)]


def set_notes(slide, text: str):
    if text and text.strip():
        slide.notes_slide.notes_text_frame.text = text.strip()


def save_image_prompt(slide_num: int, title: str, prompt: str):
    """Speichert einen Bild-Prompt als .txt im Grafiken-Ordner."""
    fname = f"{slide_num:02d}_{slugify(title)}.txt"
    full_prompt = f"{prompt}\n\n---\nStilzusatz (immer anhängen):\n{IMG_STYLE}"
    (GRAFDIR / fname).write_text(full_prompt, encoding='utf-8')
    return fname


# ── Markdown-Parser (identisch mit generate_pptx.py, eigenständig) ───────────
def parse_markdown(path: str) -> list:
    text = Path(path).read_text('utf-8')
    raw  = re.split(r'\n---\n', text)
    blocks = []
    for blk in raw:
        blk = blk.strip()
        if not blk:
            continue
        first_line = blk.split('\n')[0].strip()
        is_known = (
            first_line.startswith('## Folie') or
            first_line.startswith('## Teil') or
            first_line.startswith('# ') or
            '## Anhang' in first_line
        )
        if is_known or not blocks:
            blocks.append(blk)
        else:
            blocks[-1] = blocks[-1] + '\n\n' + blk

    slides = []
    for block in blocks:
        lines = block.split('\n')
        if not lines:
            continue
        first = lines[0].strip()

        m = re.match(r'^## Teil (\d+)[:\.]?\s*(.+)$', first)
        if m:
            slides.append({'type': 'section', 'num': int(m.group(1)),
                           'title': m.group(2).strip(), 'notes': ''})
            continue

        if first.startswith('## Anhang'):
            table_buf = []
            for line in lines:
                if line.startswith('|'):
                    if re.match(r'^\|[-| :]+\|$', line.strip()): continue
                    cells = [c.strip() for c in line.strip().strip('|').split('|')]
                    table_buf.append(cells)
            slides.append({'type': 'glossar', 'title': 'Anhang: Glossar',
                           'table_data': table_buf, 'notes': '',
                           'image_prompt': '', 'bullets': []})
            continue

        m = re.match(r'^## Folie (\d+)\s+(?:–|-)\s+(.+)$', first)
        if not m:
            continue

        num   = int(m.group(1))
        title = m.group(2).strip()
        sd = {'type': 'content', 'num': num, 'title': title,
              'bullets': [], 'table_data': [], 'image_prompt': '',
              'notes': '', 'main_text': title, 'subtitle_text': ''}

        in_notes = False
        in_img   = False
        notes_buf = []
        table_buf = []

        for line in lines[1:]:
            if line.strip() == '>':
                if in_notes: notes_buf.append('')
                continue
            if line.startswith('> '):
                content = line[2:].strip()
                if '**Sprecher-Notizen:**' in content:
                    in_notes = True; in_img = False
                elif '**Bild-Prompt:**' in content:
                    in_img = True; in_notes = False
                elif content.startswith('*Quelle:'):
                    in_notes = False; in_img = False
                elif in_img and content.startswith('`') and content.endswith('`'):
                    sd['image_prompt'] = content[1:-1]
                    in_img = False
                elif in_notes and content:
                    clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)
                    clean = re.sub(r'\*([^*]+)\*', r'\1', clean)
                    notes_buf.append(clean)
                continue
            if line.startswith('|'):
                if re.match(r'^\|[-| :]+\|$', line.strip()): continue
                cells = [c.strip() for c in line.strip().strip('|').split('|')]
                table_buf.append(cells)
                continue
            nm = re.match(r'^(\d+)\.\s+(.+)$', line)
            if nm:
                sd['bullets'].append({'text': nm.group(2).strip(), 'num': int(nm.group(1)),
                                      'section_header': False, 'indent': 0})
                continue
            if line.startswith('- '):
                sd['bullets'].append({'text': line[2:].strip(), 'num': None,
                                      'section_header': False, 'indent': 0})
                continue
            sh = re.match(r'^\*\*(.+)\*\*$', line.strip())
            if sh:
                sd['bullets'].append({'text': sh.group(1), 'num': None,
                                      'section_header': True, 'indent': 0})
                continue
            if num == 1:
                bt = re.match(r'^\*\*(.+)\*\*$', line.strip())
                if bt: sd['main_text'] = bt.group(1); continue
                it = re.match(r'^\*([^*].+)\*$', line.strip())
                if it: sd['subtitle_text'] = it.group(1); continue

        if table_buf: sd['table_data'] = table_buf
        sd['notes'] = '\n'.join(notes_buf).strip()

        if num == 1: sd['type'] = 'title'
        elif num == 36: sd['type'] = 'quote'
        elif sd['table_data'] and not sd['bullets']: sd['type'] = 'table'
        elif sd['table_data'] and sd['bullets']: sd['type'] = 'mixed'

        slides.append(sd)

    return slides


# ── Placeholder-Helfer ────────────────────────────────────────────────────────
def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_title(slide, text, idx=0):
    ph = get_ph(slide, idx)
    if ph:
        ph.text = text


def calc_fs(n: int) -> int:
    if n <= 4:  return 18
    if n <= 6:  return 16
    if n <= 9:  return 14
    return 12


def write_bullets(tf, bullets):
    """Bullets in ein TextFrame schreiben (behält Template-Formatierung weitgehend)."""
    tf.clear()
    tf.word_wrap = True
    n  = len([b for b in bullets if not b.get('section_header')])
    fs = calc_fs(n)
    first = True

    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5 if b.get('section_header') else 2)
        p.alignment    = PP_ALIGN.LEFT

        if b.get('section_header'):
            r = p.add_run()
            r.text = b['text']
            r.font.size = Pt(fs); r.font.bold = True
        else:
            num = b.get('num')
            if num is not None:
                r0 = p.add_run(); r0.text = f"{num}.  "
                r0.font.size = Pt(fs); r0.font.bold = True
            else:
                r0 = p.add_run(); r0.text = "●  "
                r0.font.size = Pt(fs); r0.font.bold = True

            for txt, bold, ital in parse_md(b['text']):
                if txt:
                    r = p.add_run(); r.text = txt
                    r.font.size = Pt(fs)
                    r.font.bold = bold; r.font.italic = ital


def add_table(slide, rows, left, top, width, height):
    if not rows: return
    nr, nc = len(rows), len(rows[0])
    rh = min(Inches(0.48), height / nr)
    tbl = slide.shapes.add_table(nr, nc, int(left), int(top),
                                  int(width), int(rh * nr)).table
    # Header
    for j, txt in enumerate(rows[0]):
        c = tbl.cell(0, j)
        c.text = ''.join(s[0] for s in parse_md(txt))
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)  # Navy
        for para in c.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(12); run.font.bold = True
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Daten
    ALT = RGBColor(0xEE, 0xF2, 0xF7)
    for i, row in enumerate(rows[1:], 1):
        for j, txt in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ''.join(s[0] for s in parse_md(txt))
            c.fill.solid()
            c.fill.fore_color.rgb = ALT if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            for para in c.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11); run.font.bold = (j == 0)
                    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)


# ── Slide-Builder ─────────────────────────────────────────────────────────────
def build(prs, sd, layout_map, grafik_log):
    t = sd.get('type', 'content')

    # Grafik-Prompt speichern (Folien mit Bild-Prompt)
    num = sd.get('num', 0)
    title = sd.get('title', '')
    prompt = sd.get('image_prompt', '')
    if prompt and num:
        fname = save_image_prompt(num, title, prompt)
        grafik_log.append(f"  Folie {num:2d}: {fname}")

    # ── Titelfolie ────────────────────────────────────────────────────────────
    if t == 'title':
        layout = layout_map.get('Titel 2', layout_map.get('Titel', prs.slide_layouts[0]))
        slide = prs.slides.add_slide(layout)
        set_title(slide, sd.get('main_text', 'Künstliche Intelligenz'))
        ph1 = get_ph(slide, 1)
        if ph1: ph1.text = sd.get('subtitle_text', '')
        set_notes(slide, sd.get('notes', ''))
        return slide

    # ── Abschnitts-Trennfolie ─────────────────────────────────────────────────
    if t == 'section':
        layout = layout_map.get('Abschnittsüberschrift 2',
                  layout_map.get('Abschnittsüberschrift', prs.slide_layouts[8]))
        slide = prs.slides.add_slide(layout)
        ph = get_ph(slide, 0)
        if ph:
            ph.text = sd['title']
            for para in ph.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
        return slide

    # ── Zitat-Folie ───────────────────────────────────────────────────────────
    if t == 'quote':
        layout = layout_map.get('Zitat 2', layout_map.get('Zitat', prs.slide_layouts[40]))
        slide = prs.slides.add_slide(layout)
        ph = get_ph(slide, 13)   # Inhaltsplatzhalter in Zitat-Layout
        if not ph: ph = get_ph(slide, 0)
        if ph:
            ph.text = (
                '„KI ist nicht die Zukunft.\n'
                'KI ist die Gegenwart.\n'
                'Die Frage ist nur:\n'
                'Wie gestalten wir sie?"'
            )
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
        set_notes(slide, sd.get('notes', ''))
        return slide

    # ── Glossar ───────────────────────────────────────────────────────────────
    if t == 'glossar':
        layout = layout_map.get('Inhalt\xa04', layout_map.get('Inhalt', prs.slide_layouts[1]))
        slide = prs.slides.add_slide(layout)
        set_title(slide, 'Anhang: Glossar')
        ph = get_ph(slide, 1)
        if ph:
            sp = ph._element; sp.getparent().remove(sp)
        rows = sd.get('table_data', [])
        if rows:
            add_table(slide, rows, Inches(0.67), Inches(2.0), Inches(11.99), Inches(4.7))
        return slide

    # ── Tabellen-Folie ────────────────────────────────────────────────────────
    if t == 'table':
        layout = layout_map.get('Nur Titel', prs.slide_layouts[45])
        slide = prs.slides.add_slide(layout)
        set_title(slide, sd['title'])
        rows = sd.get('table_data', [])
        if rows:
            add_table(slide, rows, Inches(0.67), Inches(2.1), Inches(11.99), Inches(4.6))
        if prompt:
            _add_img_note(slide, prompt)
        set_notes(slide, sd.get('notes', ''))
        return slide

    # ── Mixed (Bullets + Tabelle) ─────────────────────────────────────────────
    if t == 'mixed':
        layout = layout_map.get('Nur Titel', prs.slide_layouts[45])
        slide = prs.slides.add_slide(layout)
        set_title(slide, sd['title'])
        bullets = sd.get('bullets', [])
        if bullets:
            tb = slide.shapes.add_textbox(
                Inches(0.67), Inches(2.0), Inches(11.99), Inches(2.0))
            write_bullets(tb.text_frame, bullets)
        rows = sd.get('table_data', [])
        if rows:
            add_table(slide, rows, Inches(0.67), Inches(4.1), Inches(11.99), Inches(2.5))
        if prompt:
            _add_img_note(slide, prompt)
        set_notes(slide, sd.get('notes', ''))
        return slide

    # ── Standard-Inhaltsfolie ─────────────────────────────────────────────────
    # Layout-Auswahl: mit Bild wenn Prompt vorhanden (Bild rechts = Layout 20)
    if prompt:
        layout = layout_map.get('Inhalt mit Bild 4', prs.slide_layouts[20])
    else:
        layout = layout_map.get('Inhalt', prs.slide_layouts[1])

    slide = prs.slides.add_slide(layout)
    set_title(slide, sd['title'])

    bullets = sd.get('bullets', [])
    ph_content = get_ph(slide, 14) or get_ph(slide, 1)
    if ph_content and bullets:
        write_bullets(ph_content.text_frame, bullets)

    # Bild-Placeholder mit Prompt-Hinweis beschriften
    if prompt:
        ph_img = get_ph(slide, 13)
        if ph_img:
            _label_image_ph(ph_img, num, title)

    set_notes(slide, sd.get('notes', ''))
    return slide


def _label_image_ph(ph, num, title):
    """Bildplatzhalter mit Dateiname und Hinweis beschriften (als Textbox über dem PH)."""
    slide = ph._element.getparent().getparent()  # <p:sp> → <p:spTree> → <p:cSld>
    # Wir schreiben den Dateinamen in die Notes – PH bleibt leer für echtes Bild
    pass  # Platzhalter bleibt sauber frei – Dateiname steht in Grafiken-Ordner


def _add_img_note(slide, prompt):
    """Bei Layouts ohne Bild-PH: kleinen Hinweis-Text unten rechts."""
    tb = slide.shapes.add_textbox(
        Inches(9.0), Inches(6.55), Inches(4.0), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "→ Bild: siehe Grafiken-Ordner"
    r.font.size = Pt(8)
    r.font.italic = True
    r.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)


# ── Haupt-Funktion ────────────────────────────────────────────────────────────
def main():
    print(f"Lese:    {MD}")
    print(f"Vorlage: {TMPL}")
    slides_data = parse_markdown(str(MD))

    prs = Presentation(str(TMPL))

    # Demo-Folien sauber entfernen
    NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    for sld_id in list(prs.slides._sldIdLst):
        rId = sld_id.get(f'{{{NS_R}}}id')
        if rId:
            try: prs.part.drop_rel(rId)
            except: pass
        prs.slides._sldIdLst.remove(sld_id)

    # Layout-Map aufbauen
    layout_map = {l.name: l for l in prs.slide_layouts}

    grafik_log = []

    print(f"\n{len(slides_data)} Elemente werden verarbeitet:\n")
    for sd in slides_data:
        t = sd.get('type', 'content')
        label = sd.get('title', f"Teil {sd.get('num','?')}")
        n     = sd.get('num', '')
        num_s = f"[{n:>3}]" if n else "     "
        print(f"  {num_s} [{t:8s}]  {label[:62]}")
        build(prs, sd, layout_map, grafik_log)

    prs.save(str(OUT))
    print(f"\n✓ Gespeichert: {OUT}")
    print(f"  {len(prs.slides)} Folien")

    # Grafiken-Prompts
    if grafik_log:
        print(f"\n📁 {len(grafik_log)} Grafik-Prompts in {GRAFDIR}/:")
        for line in grafik_log:
            print(line)

    # Gesamt-Stilhinweis
    style_file = GRAFDIR / '00_STIL_FUER_ALLE_BILDER.txt'
    style_file.write_text(
        "Stilzusatz – an JEDEN Bildprompt anhängen:\n\n"
        f"{IMG_STYLE}\n\n"
        "Empfohlene Tools:\n"
        "  • Copilot Designer (Microsoft)\n"
        "  • DALL·E in ChatGPT\n"
        "  • Adobe Firefly\n"
        "  • Midjourney\n\n"
        "Nach Bildgenerierung:\n"
        "  Bild in PowerPoint einfügen: Bildplatzhalter auf Folie anklicken → Bild einfügen\n"
        "  Zieldateiname: NN_Folientitel.png (NN = Foliennummer)\n",
        encoding='utf-8'
    )
    print(f"\n  + Stilanleitung: {style_file.name}")


if __name__ == '__main__':
    main()
